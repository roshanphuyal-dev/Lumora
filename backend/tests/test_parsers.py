import io
from unittest.mock import patch

import pytest
from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfWriter

from app.parsers import UnsupportedFileTypeError, get_parser
from app.parsers.docx_parser import parse as parse_docx
from app.parsers.image_parser import parse as parse_image
from app.parsers.pdf_parser import parse as parse_pdf
from app.parsers.pptx_parser import parse as parse_pptx


def _pdf_bytes(page_count: int = 2) -> bytes:
    writer = PdfWriter()
    for _ in range(page_count):
        writer.add_blank_page(width=200, height=200)
    buf = io.BytesIO()
    writer.write(buf)
    return buf.getvalue()


def _pptx_bytes(slide_texts: list[str]) -> bytes:
    presentation = Presentation()
    layout = presentation.slide_layouts[1]
    for text in slide_texts:
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = text
    buf = io.BytesIO()
    presentation.save(buf)
    return buf.getvalue()


def _docx_bytes(paragraphs: list[str]) -> bytes:
    document = DocxDocument()
    for paragraph in paragraphs:
        document.add_paragraph(paragraph)
    buf = io.BytesIO()
    document.save(buf)
    return buf.getvalue()


def _png_bytes() -> bytes:
    from PIL import Image

    buf = io.BytesIO()
    Image.new("RGB", (10, 10), color="white").save(buf, format="PNG")
    return buf.getvalue()


def test_pdf_parser_reports_page_count_and_sections() -> None:
    result = parse_pdf(_pdf_bytes(page_count=3))

    assert result.page_count == 3
    assert [section.index for section in result.sections] == [1, 2, 3]


def test_pptx_parser_extracts_one_section_per_slide() -> None:
    result = parse_pptx(_pptx_bytes(["Slide One", "Slide Two"]))

    assert result.page_count == 2
    assert [section.text for section in result.sections] == ["Slide One", "Slide Two"]
    assert "Slide One" in result.text
    assert "Slide Two" in result.text


def test_docx_parser_joins_paragraphs_into_a_single_section() -> None:
    result = parse_docx(_docx_bytes(["First paragraph.", "Second paragraph."]))

    assert result.page_count is None
    assert len(result.sections) == 1
    assert "First paragraph." in result.text
    assert "Second paragraph." in result.text


def test_image_parser_returns_ocr_text() -> None:
    with patch("app.parsers.image_parser.pytesseract.image_to_string", return_value="hello world"):
        result = parse_image(_png_bytes())

    assert result.text == "hello world"
    assert result.page_count == 1
    assert result.sections[0].text == "hello world"


def test_image_parser_with_no_text_returns_no_sections() -> None:
    with patch("app.parsers.image_parser.pytesseract.image_to_string", return_value="   "):
        result = parse_image(_png_bytes())

    assert result.text == ""
    assert result.sections == []


@pytest.mark.parametrize(
    ("mime_type", "file_type"),
    [
        ("application/pdf", "pdf"),
        (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "pptx",
        ),
        ("application/vnd.openxmlformats-officedocument.wordprocessingml.document", "docx"),
        ("image/png", "png"),
    ],
)
def test_get_parser_dispatches_by_mime_type(mime_type: str, file_type: str) -> None:
    assert get_parser(mime_type, file_type) is not None


def test_get_parser_falls_back_to_file_type_when_mime_type_unknown() -> None:
    assert get_parser("application/octet-stream", "docx") is parse_docx


def test_get_parser_raises_for_unknown_type() -> None:
    with pytest.raises(UnsupportedFileTypeError):
        get_parser("application/x-nonsense", "xyz")
