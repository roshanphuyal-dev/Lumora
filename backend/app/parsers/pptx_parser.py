"""PPTX text extraction (python-pptx) — see app/parsers/registry.py for dispatch."""

import io

from pptx import Presentation

from app.parsers.base import ParsedDocument, ParsedSection


def parse(file_bytes: bytes) -> ParsedDocument:
    """Extract per-slide text and metadata from a PPTX's raw bytes."""
    presentation = Presentation(io.BytesIO(file_bytes))

    sections = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts = [
            shape.text_frame.text.strip()
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        sections.append(ParsedSection(index=index, text="\n".join(texts)))

    text = "\n\n".join(section.text for section in sections if section.text)
    title = presentation.core_properties.title or None

    return ParsedDocument(
        text=text, sections=sections, page_count=len(presentation.slides), title=title
    )
